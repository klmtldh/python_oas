import pptx
#from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches,Centipoints,Cm,Emu,Mm,Pt
from pptx.enum.dml import MSO_THEME_COLOR
from pptx.dml.color import ColorFormat,RGBColor
from pptx.enum.text import PP_ALIGN

from pptx.chart.data import ChartData,CategoryChartData
from pptx.enum.chart import XL_CHART_TYPE

from pptx.enum.chart import XL_LABEL_POSITION
from pptx.enum.chart import XL_LEGEND_POSITION


class OasPptx(object):
    '''
    p_slide_layout=1
    若指定了PPT文件，p_slide_layout的数字为指定PPT中第n页的布局。
    pptx的文档结构是pptx->slide->shape(text frame)->paragraph->run(一段话中的任意部分)，所以程序处理也差不多按着这个顺序来。
    Presentation pptx文档
    slide 幻灯片——layout 布局
        auto shape 图形(shapes)
            title 标题
            placeholders 占位符
                text_frame 框架
                    paragraph 段落
            textbox 文本框
        picture 图片
        graphic frame 图形框
        group shape 组合图形
        line/connector 线/连接器
        content part 内容部分


        @author: Kongling

    '''

    def __init__(self,
                 p_path=None,
                 p_slide_layout=1,
                 f_name='微软雅黑',
                 p_height=6858000, p_width=12192000
                 ):
        self.p_path = p_path
        self.p_slide_layout = p_slide_layout
        self.f_name = f_name
        # 实例化Presentation为prs
        self.prs = pptx.Presentation(p_path)
        # 设置slide幻灯片layout布局，并赋予slide_layout
        self.slide_layout = self.prs.slide_layouts[self.p_slide_layout]
        # 设置PPT为16：9的宽屏
        self.prs.slide_height = p_height  # 设置ppt的高度
        self.prs.slide_width = p_width  # 设置ppt的宽度
        # 添加页面
        self.slide = self.prs.slides.add_slide(self.slide_layout)
        # 添加形状
        self.shapes = self.slide.shapes

    def alter_csg_logo_first_page(self, text_list):
        for i, shape in enumerate(self.prs.slides[0].shapes):
            if not shape.has_text_frame:
                continue
            text_frame = shape.text_frame
            # text_frame.clear()
            p = text_frame.paragraphs[0]
            p.text = text_list[i]
            if i == 0:
                p.font.bold = True
                # 设置字体，一旦有中文就不正常，英文还好
                p.font.name = '微软雅黑'
                # 设置字体大小
                p.font.size = Pt(44)
                # 设置颜色
                p.font.color.rgb = RGBColor(0, 54, 122)
            elif i == 1:
                p.font.bold = True
                # 设置字体，一旦有中文就不正常，英文还好
                p.font.name = '微软雅黑'
                # 设置字体大小
                p.font.size = Pt(24)
                # 设置颜色
                p.font.color.rgb = RGBColor(0, 0, 0)
            elif i == 2:
                p.font.bold = True
                # 设置字体，一旦有中文就不正常，英文还好
                p.font.name = '微软雅黑'
                # 设置字体大小
                p.font.size = Pt(24)
                # 设置颜色
                p.font.color.rgb = RGBColor(0, 0, 0)
            elif i == 3:
                p.font.bold = True
                # 设置字体，一旦有中文就不正常，英文还好
                p.font.name = '微软雅黑'
                # 设置字体大小
                p.font.size = Pt(18)
                # 设置颜色
                p.font.color.rgb = RGBColor(0, 54, 122)

        # text_frame = self.prs.slides[0].shapes[0].text_frame
        # # text_frame.clear()
        # p = text_frame.paragraphs[0]
        # p.text = text_list[0]
        # p.font.bold = True
        # # 设置字体，一旦有中文就不正常，英文还好
        # p.font.name = '微软雅黑'
        # # 设置字体大小
        # p.font.size = Pt(44)
        # # 设置颜色
        # p.font.color.rgb = RGBColor(0, 54, 122)
        #
        # text_frame = self.prs.slides[0].shapes[1].text_frame
        # p = text_frame.paragraphs[0]
        # p.text = text_list[1]
        # p.font.bold = True
        # # 设置字体，一旦有中文就不正常，英文还好
        # p.font.name = '微软雅黑'
        # # 设置字体大小
        # p.font.size = Pt(24)
        # # 设置颜色
        # p.font.color.rgb = RGBColor(0, 0, 0)
        #
        # text_frame = self.prs.slides[0].shapes[2].text_frame
        # p = text_frame.paragraphs[0]
        # p.text = text_list[2]
        # p.font.bold = True
        # # 设置字体，一旦有中文就不正常，英文还好
        # p.font.name = '微软雅黑'
        # # 设置字体大小
        # p.font.size = Pt(24)
        # # 设置颜色
        # p.font.color.rgb = RGBColor(0, 0, 0)

    def page(self):
        # 添加slide幻灯片，并赋予slide
        self.slide = self.prs.slides.add_slide(self.slide_layout)
        self.shapes = self.slide.shapes

    def line(self):
        slide = 0

    def text(self, text, font_name=None, font_bold=False, alig='left', font_size=22, font_color=(0, 0, 0),
             textbox_left=2, textbox_top=2, textbox_width=90, textbox_height=20):
        textbox_left = int(0.01 * textbox_left * self.prs.slide_width)
        textbox_top = int(0.01 * textbox_top * self.prs.slide_height)
        textbox_width = int(0.01 * textbox_width * self.prs.slide_width)
        textbox_height = int(0.01 * textbox_height * self.prs.slide_height)
        if font_name == None:
            font_name = self.f_name
        # 在指定位置插入一个文本框，我按比例填的。
        tBox = self.slide.shapes.add_textbox(left=textbox_left,
                                             top=textbox_top,
                                             width=textbox_width,
                                             height=textbox_height)
        # 格式化为文本格式
        tf = tBox.text_frame
        tf.word_wrap = True
        # 插入段落
        p = tf.add_paragraph()
        # 设置对齐方式
        if alig == 'left':
            p.alignment = PP_ALIGN.LEFT
        elif alig == 'center':
            p.alignment = PP_ALIGN.CENTER

        p.level = 1
        # 设置粗体
        p.font.bold = font_bold
        # 设置字体，一旦有中文就不正常，英文还好
        p.font.name = font_name
        # 设置字体大小
        p.font.size = Pt(font_size)
        # 设置颜色
        p.font.color.rgb = RGBColor(font_color[0], font_color[1], font_color[2])
        # 设置文本内容
        p.text = text

    def text_n(self, text, font_name=None, font_bold=False, alig='left', font_size=22, font_color=(0, 0, 0),
               textbox_left=2, textbox_top=2, textbox_width=90, textbox_height=20):
        textbox_left = int(0.01 * textbox_left * self.prs.slide_width)
        textbox_top = int(0.01 * textbox_top * self.prs.slide_height)
        textbox_width = int(0.01 * textbox_width * self.prs.slide_width)
        textbox_height = int(0.01 * textbox_height * self.prs.slide_height)
        if font_name == None:
            font_name = self.f_name
        # 在指定位置插入一个文本框，我按比例填的。
        tBox = self.slide.shapes.add_textbox(left=textbox_left,
                                             top=textbox_top,
                                             width=textbox_width,
                                             height=textbox_height)
        # 格式化为文本格式
        tf = tBox.text_frame
        tf.word_wrap = True
        if isinstance(text, str):
            # 插入段落
            p = tf.add_paragraph()
            # 设置粗体
            p.font.bold = font_bold
            # 设置字体，一旦有中文就不正常，英文还好
            p.font.name = font_name
            # 设置字体大小
            p.font.size = Pt(font_size)
            # 设置颜色
            p.font.color.rgb = RGBColor(font_color[0], font_color[1], font_color[2])
            # 设置文本内容
            p.text = '       ' + text
        elif isinstance(text, list):
            for i in range(len(text)):
                # 插入段落
                p = tf.add_paragraph()
                # 设置粗体
                p.font.bold = font_bold
                # 设置字体，一旦有中文就不正常，英文还好
                p.font.name = font_name
                # 设置字体大小
                p.font.size = Pt(font_size)
                # 设置颜色
                p.font.color.rgb = RGBColor(font_color[0], font_color[1], font_color[2])
                # 设置文本内容
                p.text = '       ' + text[i]

    def pic(self, img_path, left=0, top=0, width=None, height=None):
        # 全屏插入一张图片，图片最好提前处理长宽比，因为是不锁比例拉伸。
        pic = self.slide.shapes.add_picture(img_path, left=left, top=top, width=width, height=height)
        # 将刚插入的图片至于底层
        # self.slide.shapes._spTree.insert(1, pic._element)

    # 插入表格，单位都是页面的百分比，1表示
    def p_table(self, df, left=2, top=25, width=95, height=60):
        rows = df.shape[0]+1
        cols = df.shape[1]
        left = int(0.01 * left * self.prs.slide_width)
        top = int(0.01 * top * self.prs.slide_height)
        width = int(0.01 * width * self.prs.slide_width)
        height = int(0.01 * height * self.prs.slide_height)
        table = self.slide.shapes.add_table(rows, cols, left, top, width, height).table

        # 获取表头列
        head = list(df)
        # 计算表格有效列宽,页面宽度90%计算平均值
        columns_width = int(width / len(head))
        # 调整表格宽度
        for i in range(len(head)):
            table.columns[i].width = columns_width
        for i in range(rows):
            if i == 0:
                for j in range(cols):
                    table.cell(i, j).text = str(head[j])
            else:
                for j in range(cols):
                    table.cell(i, j).text = str(df.iloc[i-1, j])

        return table

    def chart(self, df, chart_class='bar', left=2, top=25, width=95, height=60):
        # 获取第一例的列名，作为categories
        categories = df.iloc[:, 0].tolist()
        # 获取第二列滞后的数据
        series_name = df.iloc[:, 1:]
        # 将第二列后的数据，按照列转换成系列
        series = series_name.to_dict(orient='series')
        # define chart data ---------------------
        chart_data = ChartData()
        chart_data.categories = categories
        for key in series:
            chart_data.add_series(key, series[key].tolist())

        left = int(0.01 * left * self.prs.slide_width)
        top = int(0.01 * top * self.prs.slide_height)
        width = int(0.01 * width * self.prs.slide_width)
        height = int(0.01 * height * self.prs.slide_height)
        if chart_class == 'bar':
            graphic_frame = self.slide.shapes.add_chart(
                XL_CHART_TYPE.COLUMN_CLUSTERED, left, top, width, height, chart_data
            )
        elif chart_class == 'pie':
            graphic_frame = self.slide.shapes.add_chart(
                XL_CHART_TYPE.PIE, left, top, width, height, chart_data
            )
        elif chart_class == 'line':
            graphic_frame = self.slide.shapes.add_chart(
                XL_CHART_TYPE.LINE, left, top, width, height, chart_data
            )

        chart = graphic_frame.chart
        if chart_class == 'line':
            chart.series[0].smooth = False
        else:
            plot = chart.plots[0]
            plot.has_data_labels = True
            data_labels = plot.data_labels

            data_labels.font.size = Pt(13)
            data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
            data_labels.position = XL_LABEL_POSITION.OUTSIDE_END

        if chart_class == 'pie':
            data_labels.number_format = '0%'

        if chart_class in ['bar', 'line'] and df.shape[1] > 2:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.RIGHT
            chart.legend.include_in_layout = False
        else:
            chart.has_legend = True
            chart.legend.position = XL_LEGEND_POSITION.BOTTOM
            chart.legend.include_in_layout = False

    # 保存为文件
    def save_pptx(self, pptx_name):
        self.prs.save(pptx_name)  # 保存文档